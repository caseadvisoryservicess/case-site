#!/usr/bin/env python3
"""Страж раскладки pptx: читает готовый файл, рисует его геометрию в картинки и
сообщает о переполнениях и выходах за границы листа.

Нужен потому, что LibreOffice Impress в среде сборки не установлен и штатно
отрендерить pptx нечем. Шрифты Carlito и Caladea (ставятся пакетами
fonts-crosextra-carlito и fonts-crosextra-caladea) метрически совпадают с
Calibri и Cambria, поэтому переносы строк честные. Для заголовков вместо
Caladea подставляется DejaVu Serif: в Caladea нет кириллицы, в Cambria есть.

Запуск: python3 scripts/check-pptx-layout.py <файл.pptx> [каталог-картинок]"""
import sys, os, glob, io
from pptx import Presentation
from pptx.util import Emu
from PIL import Image, ImageDraw, ImageFont

DPI = 130
EMU = 914400.0

FONTS = {}
def _find(*names):
    for n in names:
        hits = glob.glob(f'/usr/share/fonts/**/{n}', recursive=True)
        if hits: return hits[0]
    return None

FF = {
    ('Calibri', False): _find('Carlito-Regular.ttf'),
    ('Calibri', True):  _find('Carlito-Bold.ttf'),
    ('Cambria', False): _find('DejaVuSerif.ttf'),
    ('Cambria', True):  _find('DejaVuSerif-Bold.ttf'),
}
FALLBACK = _find('DejaVuSans.ttf')

def font(name, size_pt, bold):
    key = (name if (name, bold) in FF else 'Calibri', bold)
    p = FF.get(key) or FALLBACK
    px = max(6, int(round(size_pt * DPI / 72.0)))
    k = (p, px)
    if k not in FONTS:
        FONTS[k] = ImageFont.truetype(p, px)
    return FONTS[k]

def px(emu): return emu / EMU * DPI

def rgb(c, default=None):
    try:
        if c and c.type is not None and str(c.type).startswith('MSO_THEME') is False:
            return '#%02X%02X%02X' % tuple(c.rgb)
    except Exception:
        pass
    return default

def shape_fill(sh):
    try:
        f = sh.fill
        if f.type is not None and int(f.type) == 1:  # solid
            return rgb(f.fore_color, None)
    except Exception:
        pass
    return None

def shape_line(sh):
    try:
        c = rgb(sh.line.color, None)
        w = sh.line.width
        return c, (px(w) if w else 1)
    except Exception:
        return None, 1

def wrap(draw, text, fnt, maxw):
    out = []
    for para in text.split('\n'):
        words, cur = para.split(' '), ''
        for w in words:
            t = (cur + ' ' + w).strip()
            if draw.textlength(t, font=fnt) <= maxw or not cur:
                cur = t
            else:
                out.append(cur); cur = w
        out.append(cur)
    return out

ISSUES = []

def draw_tf(draw, tf, x, y, w, h, sid, label):
    """Рисует текстовый фрейм и проверяет, влезает ли он."""
    pad = 0.05 * DPI
    ax, ay, aw = x + pad, y + pad, max(10, w - 2 * pad)
    lines = []
    for p in tf.paragraphs:
        runs = [r for r in p.runs if r.text]
        if not runs:
            lines.append((None, '', 0)); continue
        r0 = runs[0]
        sz = (r0.font.size.pt if r0.font.size else 14)
        bold = bool(r0.font.bold)
        name = r0.font.name or 'Calibri'
        col = rgb(r0.font.color, '#212326') or '#212326'
        fnt = font(name, sz, bold)
        txt = ''.join(r.text for r in runs)
        bullet = b'buChar' in p._p.xml.encode() or b'buAutoNum' in p._p.xml.encode()
        if bullet: txt = '•  ' + txt
        for ln in wrap(draw, txt, fnt, aw):
            lines.append((fnt, ln, sz))
        lines.append((None, '', sz * 0.28))
        align = str(p.alignment or '')
        for i in range(len(lines)):
            pass
        lines[-1] = (None, '', sz * 0.28)
        if lines:
            lines[-1] = (None, '', sz * 0.28)
        # запомним выравнивание последнего блока
        for i in range(len(lines)):
            pass
        setattr(draw, '_align', align)
    total = 0
    for fnt, ln, sz in lines:
        total += (fnt.size * 1.22) if fnt else (sz * DPI / 72.0 * 0.28)
    # вертикальное выравнивание
    va = str(tf.vertical_anchor or '')
    oy = ay
    if 'MIDDLE' in va and total < h - 2 * pad:
        oy = y + (h - total) / 2
    cy = oy
    for fnt, ln, sz in lines:
        if fnt is None:
            cy += sz * DPI / 72.0 * 0.28; continue
        # выравнивание берём из первого параграфа
        tw = draw.textlength(ln, font=fnt)
        cx = ax
        al = None
        for p in tf.paragraphs:
            if p.runs and ''.join(r.text for r in p.runs) and ln.strip() and ln.strip() in ''.join(r.text for r in p.runs):
                al = str(p.alignment or ''); break
        if al and 'CENTER' in al: cx = x + (w - tw) / 2
        elif al and 'RIGHT' in al: cx = x + w - pad - tw
        col = '#212326'
        for p in tf.paragraphs:
            for r in p.runs:
                if ln.replace('•  ', '')[:12] and ln.replace('•  ', '')[:12] in r.text:
                    col = rgb(r.font.color, '#212326') or '#212326'; break
        draw.text((cx, cy), ln, font=fnt, fill=col)
        cy += fnt.size * 1.22
    if total > h - pad and total > 0:
        ISSUES.append(f'слайд {sid}: текст выше блока на {round((total-h+pad)/DPI,2)}" [{label}] "{(lines[0][1] if lines else "")[:46]}"')
    return total

def main(src, outdir):
    prs = Presentation(src)
    SW, SH = px(prs.slide_width), px(prs.slide_height)
    os.makedirs(outdir, exist_ok=True)
    for idx, slide in enumerate(prs.slides, 1):
        bg = '#FFFFFF'
        xml = slide._element.xml
        if 'bg>' in xml:
            import re
            m = re.search(r'<p:bg>.*?srgbClr val="([0-9A-Fa-f]{6})"', xml, re.S)
            if m: bg = '#' + m.group(1)
        im = Image.new('RGB', (int(SW), int(SH)), bg)
        d = ImageDraw.Draw(im)
        for sh in slide.shapes:
            try:
                x, y, w, h = px(sh.left), px(sh.top), px(sh.width), px(sh.height)
            except Exception:
                continue
            if x < -2 or y < -2 or x + w > SW + 2 or y + h > SH + 2:
                ISSUES.append(f'слайд {idx}: фигура выходит за лист ({round(x/DPI,2)},{round(y/DPI,2)} {round(w/DPI,2)}x{round(h/DPI,2)}) {sh.shape_type}')
            st = str(sh.shape_type)
            if 'PICTURE' in st:
                try:
                    pic = Image.open(io.BytesIO(sh.image.blob)).convert('RGB')
                    tw, th = int(w), int(h)
                    sc = max(tw / pic.width, th / pic.height)
                    pic = pic.resize((max(1, int(pic.width * sc)), max(1, int(pic.height * sc))), Image.LANCZOS)
                    ox, oy = (pic.width - tw) // 2, (pic.height - th) // 2
                    im.paste(pic.crop((ox, oy, ox + tw, oy + th)), (int(x), int(y)))
                except Exception as e:
                    d.rectangle([x, y, x + w, y + h], outline='#FF0000')
                continue
            if 'TABLE' in st or getattr(sh, 'has_table', False):
                tbl = sh.table
                cy = y
                nrows = len(tbl.rows)
                for ri, row in enumerate(tbl.rows):
                    rh = px(row.height)
                    cx = x
                    for ci, cell in enumerate(row.cells):
                        cw = px(tbl.columns[ci].width)
                        fillc = None
                        try:
                            if cell.fill.type is not None and int(cell.fill.type) == 1:
                                fillc = rgb(cell.fill.fore_color, None)
                        except Exception: pass
                        d.rectangle([cx, cy, cx + cw, cy + rh], fill=fillc, outline='#DFDCD5')
                        draw_tf(d, cell.text_frame, cx, cy, cw, rh, idx, f'таблица r{ri}c{ci}')
                        cx += cw
                    cy += rh
                if cy > SH - 0.15 * DPI:
                    ISSUES.append(f'слайд {idx}: таблица уходит вниз до {round(cy/DPI,2)}" при листе {round(SH/DPI,2)}"')
                continue
            fl = shape_fill(sh)
            lc, lw = shape_line(sh)
            if fl or lc:
                if 'OVAL' in st or 'ELLIPSE' in st:
                    d.ellipse([x, y, x + w, y + h], fill=fl, outline=lc, width=max(1, int(lw)))
                else:
                    d.rectangle([x, y, x + w, y + h], fill=fl, outline=lc, width=max(1, int(lw)))
            if sh.has_text_frame and sh.text_frame.text.strip():
                draw_tf(d, sh.text_frame, x, y, w, h, idx, st)
        im.save(os.path.join(outdir, f'slide-{idx:02d}.jpg'), quality=88)
    print(f'слайдов: {len(prs.slides)} -> {outdir}')
    if ISSUES:
        print('\nПОДОЗРЕНИЯ (%d):' % len(ISSUES))
        for i in ISSUES: print(' -', i)
    else:
        print('переполнений и выходов за лист не найдено')

if __name__ == '__main__':
    main(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else 'render')
